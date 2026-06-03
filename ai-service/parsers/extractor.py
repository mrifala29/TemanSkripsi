"""
Document text extraction — PDF and PPTX.

Responsibilities:
  - Accept raw file bytes + file type string
  - Return plain text suitable for chunking
  - Raise ValueError for unsupported types or empty extraction
"""
from __future__ import annotations

import io

import fitz  # PyMuPDF
from pptx import Presentation


def extract_text(data: bytes, file_type: str) -> str:
    """
    Extract plain text from a document binary.

    Args:
        data:      Raw file bytes downloaded from Supabase Storage.
        file_type: One of "pdf", "ppt", "pptx" (case-insensitive).

    Returns:
        Extracted text as a single string.

    Raises:
        ValueError: For unsupported file types or empty documents.
    """
    ext = file_type.lower().lstrip(".")

    if ext == "pdf":
        text = _from_pdf(data)
    elif ext in ("ppt", "pptx"):
        text = _from_pptx(data)
    else:
        raise ValueError(
            f"Unsupported file type '{file_type}'. Supported: pdf, pptx"
        )

    if not text.strip():
        raise ValueError("No text could be extracted from the document")

    return text


# ──────────────────────────────────────────────────────────────────────────
# Private helpers
# ──────────────────────────────────────────────────────────────────────────


def extract_pages_from_pdf(data: bytes) -> list[str]:
    """
    Extract text per page from a PDF binary.

    Returns:
        List of strings, one per page (1-indexed by position).
        Empty pages are kept as empty strings to preserve page numbering.

    Used by KesamaanAgent for page-accurate typo detection.
    """
    doc = fitz.open(stream=data, filetype="pdf")
    return [page.get_text() for page in doc]


def extract_pages_from_text(text: str, approx_lines_per_page: int = 40) -> list[str]:
    """
    Split extracted text into logical "pages" for typo detection.

    Since raw extracted text doesn't have page boundaries, we approximate
    by splitting every N lines (default 40 lines ≈ 1 physical page).

    Args:
        text: Full extracted text from document
        approx_lines_per_page: Lines per logical page (adjust for document density)

    Returns:
        List of text chunks representing "pages"
    """
    lines = text.split('\n')
    pages = []
    
    for i in range(0, len(lines), approx_lines_per_page):
        page_text = '\n'.join(lines[i:i + approx_lines_per_page])
        if page_text.strip():
            pages.append(page_text)
    
    return pages if pages else [text]  # Return at least 1 page if text not empty


def _from_pdf(data: bytes) -> str:
    doc = fitz.open(stream=data, filetype="pdf")
    return "\n\n".join(page.get_text() for page in doc)


def _from_pptx(data: bytes) -> str:
    prs = Presentation(io.BytesIO(data))
    slides: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        texts = [
            shape.text_frame.text
            for shape in slide.shapes
            if shape.has_text_frame
        ]
        slide_text = "\n".join(t.strip() for t in texts if t.strip())
        if slide_text:
            slides.append(f"[Slide {i}]\n{slide_text}")
    return "\n\n".join(slides)
