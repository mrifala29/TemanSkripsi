"""
TemanSkripsi AI Service
FastAPI service for document parsing, embedding, simulation, analysis and similarity.
"""

import io
import os
import textwrap
import logging
from contextlib import asynccontextmanager
from typing import Any

import fitz  # PyMuPDF
import google.generativeai as genai
from dotenv import load_dotenv
from fastapi import BackgroundTasks, FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from pptx import Presentation
from pydantic import BaseModel
from supabase import Client, create_client

load_dotenv()

# ---------------------------------------------------------------------------
# Logging
# ---------------------------------------------------------------------------
logging.basicConfig(level=logging.INFO, format="%(asctime)s %(levelname)s %(message)s")
log = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Config
# ---------------------------------------------------------------------------
SUPABASE_URL: str = os.environ["SUPABASE_URL"]
SUPABASE_KEY: str = os.environ["SUPABASE_SERVICE_ROLE_KEY"]
GEMINI_API_KEY: str = os.environ["GEMINI_API_KEY"]

CHUNK_SIZE = 800       # characters per chunk
CHUNK_OVERLAP = 100    # characters of overlap between chunks
EMBEDDING_MODEL = "models/text-embedding-004"
CHAT_MODEL = "gemini-1.5-flash"

# Schema-mapped aspect codes
ASPECT_CODE_MAP = {
    "Kejelasan Latar Belakang":         "latar_belakang",
    "Rumusan Masalah & Tujuan":         "rumusan_masalah",
    "Kekuatan Metodologi":              "metodologi",
    "Kualitas Analisis & Pembahasan":   "analisis_pembahasan",
    "Konsistensi Pembahasan":           "konsistensi",
    "Kualitas Kesimpulan":              "kesimpulan",
}

# Evaluation aspects (must match docs/evaluation-criteria.md)
EVAL_ASPECTS = [
    "Kejelasan Latar Belakang",
    "Rumusan Masalah & Tujuan",
    "Kekuatan Metodologi",
    "Kualitas Analisis & Pembahasan",
    "Konsistensi Pembahasan",
    "Kualitas Kesimpulan",
]


# ---------------------------------------------------------------------------
# Supabase + Gemini clients (initialised at startup)
# ---------------------------------------------------------------------------
supabase: Client
embed_model: Any


@asynccontextmanager
async def lifespan(app: FastAPI):
    global supabase, embed_model
    supabase = create_client(SUPABASE_URL, SUPABASE_KEY)
    genai.configure(api_key=GEMINI_API_KEY)
    log.info("AI service ready — Supabase + Gemini configured")
    yield


app = FastAPI(title="TemanSkripsi AI Service", version="1.0.0", lifespan=lifespan)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:8000", "http://backend:8000"],
    allow_methods=["*"],
    allow_headers=["*"],
)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _chunk_text(text: str) -> list[str]:
    """Split text into overlapping chunks."""
    chunks: list[str] = []
    start = 0
    while start < len(text):
        end = start + CHUNK_SIZE
        chunks.append(text[start:end])
        start += CHUNK_SIZE - CHUNK_OVERLAP
    return [c.strip() for c in chunks if c.strip()]


def _extract_text_from_pdf(data: bytes) -> str:
    doc = fitz.open(stream=data, filetype="pdf")
    return "\n\n".join(page.get_text() for page in doc)


def _extract_text_from_pptx(data: bytes) -> str:
    prs = Presentation(io.BytesIO(data))
    slides: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        texts = [shape.text_frame.text for shape in slide.shapes if shape.has_text_frame]
        slide_text = "\n".join(t for t in texts if t.strip())
        if slide_text:
            slides.append(f"[Slide {i}]\n{slide_text}")
    return "\n\n".join(slides)


def _get_embedding(text: str) -> list[float]:
    result = genai.embed_content(model=EMBEDDING_MODEL, content=text)
    return result["embedding"]


def _parse_and_embed_document(document_id: str, file_path: str, file_type: str):
    """Background task: download → parse → chunk → embed → store."""
    log.info(f"[parse] Starting: {document_id} ({file_path})")
    try:
        # 1. Download file from Supabase Storage
        data: bytes = supabase.storage.from_("documents").download(file_path)

        # 2. Extract text
        ext = file_type.lower().lstrip(".")
        if ext == "pdf":
            text = _extract_text_from_pdf(data)
        elif ext in ("ppt", "pptx"):
            text = _extract_text_from_pptx(data)
        else:
            raise ValueError(f"Unsupported file type: {file_type}")

        if not text.strip():
            raise ValueError("No text extracted from document")

        log.info(f"[parse] Extracted {len(text)} chars from {document_id}")

        # 3. Chunk
        chunks = _chunk_text(text)
        log.info(f"[parse] {len(chunks)} chunks created")

        # 4. Embed + store chunks (batch to avoid rate limits)
        rows = []
        for i, chunk in enumerate(chunks):
            embedding = _get_embedding(chunk)
            rows.append({
                "document_id":  document_id,
                "chunk_index":  i,
                "content":      chunk,
                "embedding":    embedding,
                "metadata":     {"char_count": len(chunk)},
            })

        # Insert in batches of 20
        batch_size = 20
        for start in range(0, len(rows), batch_size):
            batch = rows[start:start + batch_size]
            supabase.table("document_chunks").insert(batch).execute()

        # 5. Mark document as done
        supabase.table("documents").update({"parse_status": "done"}).eq("id", document_id).execute()
        log.info(f"[parse] Done: {document_id}")

    except Exception as exc:
        log.error(f"[parse] Failed {document_id}: {exc}", exc_info=True)
        supabase.table("documents").update({"parse_status": "failed"}).eq("id", document_id).execute()


def _retrieve_context(document_id: str, query: str, top_k: int = 5) -> str:
    """RAG: embed query → find top-k similar chunks via pgvector RPC."""
    query_embedding = _get_embedding(query)
    result = supabase.rpc("match_document_chunks", {
        "query_embedding": query_embedding,
        "match_document_id": document_id,
        "match_count": top_k,
    }).execute()
    chunks = result.data or []
    return "\n\n---\n\n".join(c["content"] for c in chunks)


# ---------------------------------------------------------------------------
# Request / Response models
# ---------------------------------------------------------------------------

class ParseRequest(BaseModel):
    document_id: str
    file_path: str
    file_type: str


class SessionStartRequest(BaseModel):
    session_id: str
    document_id: str


class SessionMessageRequest(BaseModel):
    session_id: str
    document_id: str
    message: str
    history: list[dict] = []


class AnalyzeRequest(BaseModel):
    document_id: str
    analysis_id: str


class SimilarityRequest(BaseModel):
    document_id: str
    similarity_id: str


# ---------------------------------------------------------------------------
# Routes
# ---------------------------------------------------------------------------

@app.get("/health")
def health():
    return {"status": "ok", "service": "TemanSkripsi AI"}


@app.post("/documents/parse", status_code=202)
def parse_document(req: ParseRequest, bg: BackgroundTasks):
    """Trigger async document parsing + embedding."""
    bg.add_task(_parse_and_embed_document, req.document_id, req.file_path, req.file_type)
    return {"message": "Parsing dimulai", "document_id": req.document_id}


@app.post("/sessions/start")
def session_start(req: SessionStartRequest):
    """Generate the first question for a simulation session."""
    context = _retrieve_context(req.document_id, "topik utama penelitian metodologi")

    prompt = textwrap.dedent(f"""
        Kamu adalah dosen penguji sidang skripsi yang tegas namun membangun.
        Berdasarkan isi skripsi berikut, buat 1 pertanyaan pembuka sidang yang relevan.
        Pertanyaan harus spesifik, tidak terlalu mudah, dan mendorong mahasiswa untuk menjelaskan pemikirannya.

        Isi Skripsi:
        {context}

        Balas HANYA dengan pertanyaan (tanpa preambul, tanpa penomoran).
    """).strip()

    model = genai.GenerativeModel(CHAT_MODEL)
    response = model.generate_content(prompt)
    question = response.text.strip()

    # Simpan pesan AI ke DB
    supabase.table("messages").insert({
        "session_id": req.session_id,
        "role":       "ai",
        "content":    question,
        "turn_index": 0,
        "is_followup": False,
    }).execute()

    return {"message": question, "role": "ai"}


@app.post("/sessions/message")
def session_message(req: SessionMessageRequest):
    """Process a user message and return the AI's next question/response."""
    # RAG: retrieve relevant chunks based on user's answer
    context = _retrieve_context(req.document_id, req.message, top_k=4)

    # Build conversation history string
    history_str = ""
    for turn in req.history[-6:]:  # last 6 turns only
        role = "Dosen" if turn.get("role") == "ai" else "Mahasiswa"
        history_str += f"{role}: {turn.get('content', '')}\n"

    prompt = textwrap.dedent(f"""
        Kamu adalah dosen penguji sidang skripsi yang tegas namun membangun.

        Percakapan sebelumnya:
        {history_str}

        Konteks dari skripsi yang relevan:
        {context}

        Jawaban terakhir mahasiswa: "{req.message}"

        Berikan respons sebagai dosen penguji:
        - Jika jawaban kurang tepat atau kurang lengkap, minta klarifikasi atau beri pertanyaan lanjutan.
        - Jika jawaban baik, beri apresiasi singkat lalu ajukan pertanyaan baru yang lebih dalam.
        - Tetap fokus pada isi skripsi.

        Balas HANYA dengan respons dosen (tanpa preambul).
    """).strip()

    model = genai.GenerativeModel(CHAT_MODEL)
    response = model.generate_content(prompt)
    reply = response.text.strip()

    return {"message": reply, "role": "ai"}


@app.post("/documents/analyze")
def analyze_document(req: AnalyzeRequest):
    """Score document on 6 evaluation aspects."""
    # Retrieve broad context for each aspect
    all_text_chunks = supabase.table("document_chunks") \
        .select("content") \
        .eq("document_id", req.document_id) \
        .limit(20) \
        .execute()

    context = "\n\n---\n\n".join(c["content"] for c in (all_text_chunks.data or []))

    prompt = textwrap.dedent(f"""
        Kamu adalah dosen penguji yang mengevaluasi skripsi.
        Berikan skor 0–100 dan umpan balik singkat untuk setiap aspek berikut berdasarkan isi skripsi.

        Aspek yang dinilai:
        1. Kejelasan Latar Belakang
        2. Rumusan Masalah & Tujuan
        3. Kekuatan Metodologi
        4. Kualitas Analisis & Pembahasan
        5. Konsistensi Pembahasan
        6. Kualitas Kesimpulan

        Isi Skripsi (sampel):
        {context[:4000]}

        Balas dalam format JSON:
        {{
          "overall": <rata-rata skor>,
          "summary": "<ringkasan 1-2 kalimat>",
          "scores": [
            {{"aspect": "Kejelasan Latar Belakang", "score": <0-100>, "feedback": "<umpan balik>"}},
            ...
          ]
        }}
    """).strip()

    model = genai.GenerativeModel(CHAT_MODEL)
    response = model.generate_content(prompt)

    # Extract JSON from response
    import json, re
    text = response.text.strip()
    match = re.search(r"\{.*\}", text, re.DOTALL)
    if not match:
        raise HTTPException(500, "Gagal mem-parse respons analisa dari AI")

    result = json.loads(match.group())

    # Store scores in DB
    scores_rows = []
    for s in result.get("scores", []):
        aspect_display = s["aspect"]
        aspect_code = ASPECT_CODE_MAP.get(aspect_display, "latar_belakang")
        scores_rows.append({
            "analysis_id": req.analysis_id,
            "aspect":      aspect_code,
            "score":       s["score"],
            "notes":       s.get("feedback", ""),
        })
    if scores_rows:
        supabase.table("analysis_scores").insert(scores_rows).execute()

    # Update analysis record
    supabase.table("analyses").update({
        "overall_score": result.get("overall", 0),
        "summary":       result.get("summary", ""),
        "status":        "done",
    }).eq("id", req.analysis_id).execute()

    return result


@app.post("/documents/similarity")
def check_similarity(req: SimilarityRequest):
    """Check document similarity against corpus using pgvector."""
    # Get all chunks for the document
    chunks_result = supabase.table("document_chunks") \
        .select("content,embedding") \
        .eq("document_id", req.document_id) \
        .limit(10) \
        .execute()

    chunks = chunks_result.data or []
    if not chunks:
        raise HTTPException(404, "Dokumen belum diproses")

    # For each chunk, find similar chunks from OTHER documents
    similarity_scores: list[float] = []
    for chunk in chunks:
        result = supabase.rpc("match_other_document_chunks", {
            "query_embedding":    chunk["embedding"],
            "exclude_document_id": req.document_id,
            "match_count":        3,
        }).execute()
        matches = result.data or []
        if matches:
            # similarity = 1 - cosine_distance (Supabase returns distance, not similarity)
            similarity_scores.extend([1 - m.get("distance", 1) for m in matches])

    overall_pct = round((sum(similarity_scores) / len(similarity_scores)) * 100, 1) if similarity_scores else 0.0

    # Update similarity record (use ai_text_percent as overall proxy)
    supabase.table("similarity_checks").update({
        "ai_text_percent": overall_pct,
        "similarity_note": f"Similarity rata-rata dari {len(chunks)} chunk: {overall_pct}%",
        "status":          "done",
    }).eq("id", req.similarity_id).execute()

    return {"similarity_percentage": overall_pct, "chunks_checked": len(chunks)}
